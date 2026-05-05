"""
Extract slide text and speaker notes from a .pptx file.

Usage:
    python extract_notes.py ../class1/slides1.pptx
    python extract_notes.py ../class1/slides1.pptx --notes-only
    python extract_notes.py ../class1/slides1.pptx --output notes.txt
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation


def extract(pptx_path: Path, notes_only: bool = False) -> str:
    prs = Presentation(pptx_path)
    out: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== Slide {i} ===")

        if not notes_only:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            out.append(text)
            out.append("")
            out.append("--- Speaker Notes ---")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            out.append(notes if notes else "(no notes)")
        else:
            out.append("(no notes slide)")

        out.append("")

    return "\n".join(out)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument("--notes-only", action="store_true", help="Skip slide body text; output only speaker notes")
    parser.add_argument("--output", "-o", type=Path, help="Write to this file instead of stdout")
    args = parser.parse_args()

    if not args.pptx.exists():
        print(f"Error: {args.pptx} not found", file=sys.stderr)
        return 1

    result = extract(args.pptx, notes_only=args.notes_only)

    if args.output:
        args.output.write_text(result, encoding="utf-8")
        print(f"Wrote {args.output}")
    else:
        print(result)

    return 0


if __name__ == "__main__":
    sys.exit(main())
